from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from doc2md.processors.base import BaseProcessor
from doc2md.models import ConversionResult
from doc2md.config import config

class PptxProcessor(BaseProcessor):
    SUPPORTED_EXTENSIONS = [".pptx", ".ppt"]

    def process(self, file_path: str) -> ConversionResult:
        prs = Presentation(file_path)
        filename = Path(file_path).stem
        self.assets = []

        lines = []

        # Extract all images first
        for rel in prs.part.rels.values():
            if "image" in rel.target_ref:
                image_data = rel.target_part.blob
                image_name = Path(rel.target_ref).name
                self.save_asset(filename, image_name, image_data)

        # Process each slide
        for slide_num, slide in enumerate(prs.slides, 1):
            lines.append(f"\n\n## Slide {slide_num}\n\n")

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            if para.level == 0:
                                lines.append(f"### {text}\n\n")
                            else:
                                lines.append(f"{text}\n\n")

                if hasattr(shape, "image"):
                    pass  # Images already extracted above

        markdown = "".join(lines)

        # Save markdown
        output_dir = Path(config.PROJECTS_DIR) / self.project_id
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / f"{filename}.md"
        output_path.write_text(markdown)

        return ConversionResult(
            markdown=markdown,
            assets=self.assets,
            project_id=self.project_id,
            filename=f"{filename}.md"
        )